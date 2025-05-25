import chromadb
import streamlit as st  # ✅ Replace dotenv with Streamlit's secrets
from langchain_openai import ChatOpenAI
from pptx import Presentation
from pptx.util import Pt

client = chromadb.PersistentClient(path="./chroma_db")
collection = client.get_collection(name="shared_memory")

# Initialize GPT for generating PPT outlines using Streamlit secrets
llm = ChatOpenAI(api_key=st.secrets["OPENAI_API_KEY"], model="gpt-4o")

def organize_ppt(topic: str, num_slides: int) -> str:
    # Get the stored summary for this topic only
    result = collection.get(ids=[f"summary_{topic}"])
    if not result["documents"]:
        return "No summary found for this topic."

    summary = result["documents"][0]

    # Use GPT to generate the outline for the PPT
    prompt = f"""
You are a professional PowerPoint presentation designer.

Using the following **topic** and **summary**, create a detailed PowerPoint outline with {num_slides} slides:
- Slide 1: Introduction (hook, context, purpose)
- Slide 2: Visual/Picture slide (suggest images)
- Slides 3-{num_slides-2}: Key insights with bullet points
- Last slide: Conclusion (summary & next steps)

**Topic:** {topic}
**Summary:**
{summary}

Generate the slides using presentation-appropriate bullet points (not paragraphs).
"""
    response = llm.invoke(prompt)
    return response.content.strip()

def generate_pptx(topic: str, num_slides: int, filename: str = "generated_presentation.pptx") -> str:
    # Fetch summary
    result = collection.get(ids=[f"summary_{topic}"])
    if not result["documents"]:
        return None

    summary = result["documents"][0]

    # Get structured slides outline from GPT
    prompt = f"""
Using the following **topic** and **summary**, create a detailed PowerPoint outline with {num_slides} slides.
Each slide should have a title and 3–5 bullet points.

**Topic:** {topic}
**Summary:** {summary}
"""
    response = llm.invoke(prompt)
    slides = response.content.strip().split("\n\n")

    # Build PPTX
    prs = Presentation()
    for slide_text in slides:
        lines = slide_text.strip().split("\n")
        if not lines:
            continue

        # Title
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = lines[0].replace("•", "").replace("-", "").strip()

        # Bullet points
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        for line in lines[1:]:
            if line.strip():
                p = tf.add_paragraph()
                p.text = line.strip().replace("-", "").replace("•", "").strip()
                p.level = 0
                p.font.size = Pt(18)

    prs.save(filename)
    return filename
